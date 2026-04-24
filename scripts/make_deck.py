"""Generate Stallmate investor deck as a .pptx file.

Run: python3 scripts/make_deck.py
Output: public/stallmate-deck.pptx (served at /stallmate/stallmate-deck.pptx)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# Colors — match the product UI
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DARK = RGBColor(0x31, 0x27, 0xA6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
GRAY_900 = RGBColor(0x11, 0x18, 0x27)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_50 = RGBColor(0xF9, 0xFA, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, left, top, width, height, fill, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=GRAY_900, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_footer(slide, slide_num, total=10):
    # Bottom accent bar
    add_rect(slide, Inches(0), Inches(7.3), SW, Inches(0.2), INDIGO)
    # Corner label
    add_text(slide, Inches(0.4), Inches(6.95), Inches(4), Inches(0.3),
             "Stallmate · Investor Deck", size=10, color=GRAY_500)
    add_text(slide, Inches(11.5), Inches(6.95), Inches(1.5), Inches(0.3),
             f"{slide_num} / {total}", size=10, color=GRAY_500, align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, subtitle=None, top=Inches(0.6)):
    if eyebrow:
        add_text(slide, Inches(0.6), top, Inches(12), Inches(0.3),
                 eyebrow.upper(), size=12, bold=True, color=INDIGO)
        top = top + Inches(0.4)
    add_text(slide, Inches(0.6), top, Inches(12), Inches(1),
             title, size=36, bold=True, color=GRAY_900)
    if subtitle:
        add_text(slide, Inches(0.6), top + Inches(0.85), Inches(12), Inches(0.6),
                 subtitle, size=18, color=GRAY_500)
    return top + Inches(1.5 if subtitle else 1.1)


# ─── SLIDE 1 — Title ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
# Indigo hero block on left
add_rect(s, 0, 0, Inches(5.5), SH, INDIGO)
# Logo circle
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(0.8),
                            Inches(0.9), Inches(0.9))
circle.fill.solid(); circle.fill.fore_color.rgb = WHITE; circle.line.fill.background()
add_text(s, Inches(0.8), Inches(0.8), Inches(0.9), Inches(0.9),
         "S", size=44, bold=True, color=INDIGO, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.9), Inches(0.95), Inches(3.5), Inches(0.7),
         "stallmate", size=28, bold=True, color=WHITE)

add_text(s, Inches(0.8), Inches(2.8), Inches(4.5), Inches(1.2),
         "Stallmate", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.1), Inches(4.5), Inches(1.6),
         "The operating system for India's stall-booking economy.",
         size=22, color=WHITE)

add_text(s, Inches(0.8), Inches(6.2), Inches(4.5), Inches(0.4),
         "SEED ROUND · 2026", size=12, bold=True, color=AMBER)

# Right side — quick visual
add_text(s, Inches(6.5), Inches(2.5), Inches(6.5), Inches(0.8),
         "50,000", size=80, bold=True, color=INDIGO)
add_text(s, Inches(6.5), Inches(3.9), Inches(6.5), Inches(0.6),
         "community events / year in India", size=20, color=GRAY_700)
add_text(s, Inches(6.5), Inches(4.6), Inches(6.5), Inches(0.6),
         "run on WhatsApp, Google Sheets, and cash.", size=18, color=GRAY_500)

add_text(s, Inches(6.5), Inches(5.6), Inches(6.5), Inches(1),
         "We're building the platform that runs them.",
         size=22, bold=True, color=GRAY_900)

add_footer(s, 1)


# ─── SLIDE 2 — Problem ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "The Problem",
            "Three sides. Three failures. Zero tooling.",
            "₹6,400 Cr of stall-booking GMV runs on WhatsApp forwards and cash.")

# Three columns
col_w = Inches(3.9)
col_top = Inches(2.7)
col_h = Inches(3.8)
cols = [
    ("🏘️  Venue Admin", "RWA / Gated Community",
     ["No idea what residents actually want",
      "Same repeat vendors, no new blood",
      "Revenue invisible — 'informal'",
      "No data to share with organisers"]),
    ("🎯  Event Manager", "Organiser",
     ["Calls 30 RWAs to find a venue",
      "Re-creates every event from scratch",
      "Payment chasing eats 20% of fees",
      "No tool to pitch vendors"]),
    ("🛍️  Vendor", "Small business / MSME",
     ["Pays ₹3K–20K upfront, no footfall signal",
      "Shows up — residents wanted plants, not sarees",
      "No financing → weaker vendors locked out",
      "No portable reputation across venues"]),
]
x = Inches(0.6)
for label, subtitle, pains in cols:
    add_rect(s, x, col_top, col_w, col_h, GRAY_50, GRAY_300)
    add_rect(s, x, col_top, col_w, Inches(0.15), INDIGO)
    add_text(s, x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6),
             Inches(0.5), label, size=20, bold=True, color=GRAY_900)
    add_text(s, x + Inches(0.3), col_top + Inches(0.75), col_w - Inches(0.6),
             Inches(0.3), subtitle, size=12, color=INDIGO, bold=True)
    y = col_top + Inches(1.3)
    for pain in pains:
        add_text(s, x + Inches(0.3), y, col_w - Inches(0.6), Inches(0.6),
                 "•  " + pain, size=14, color=GRAY_700)
        y += Inches(0.6)
    x += col_w + Inches(0.15)

add_footer(s, 2)


# ─── SLIDE 3 — Solution ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "The Solution",
            "One platform. Three dashboards. Five data loops.",
            "Stallmate closes the feedback loop no one else has.")

# Three-box flow diagram
box_top = Inches(3)
box_h = Inches(2.8)
box_w = Inches(3.6)
labels = [
    ("Venue Admin", "Owns Aparna Sarovar.\nShares demand poll in\nresident WhatsApp group.", INDIGO),
    ("Event Manager", "Sees 'residents want\nfood 40%, fashion 25%'\npre-filled in creation.", AMBER),
    ("Vendor", "Books stall with real\ndemand signal.\nGets BNPL + reputation.", GREEN),
]
x = Inches(0.6)
arrow_positions = []
for name, body, color in labels:
    add_rect(s, x, box_top, box_w, box_h, WHITE, GRAY_300)
    add_rect(s, x, box_top, box_w, Inches(0.35), color)
    add_text(s, x + Inches(0.3), box_top + Inches(0.55), box_w - Inches(0.6),
             Inches(0.5), name, size=22, bold=True, color=GRAY_900)
    add_text(s, x + Inches(0.3), box_top + Inches(1.2), box_w - Inches(0.6),
             Inches(1.5), body, size=14, color=GRAY_700)
    arrow_positions.append(x + box_w)
    x += box_w + Inches(0.55)

# Draw arrows between boxes
for ax in arrow_positions[:-1]:
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.05),
                               box_top + Inches(1.2), Inches(0.45), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = GRAY_300
    arrow.line.fill.background()

# Feedback loop arrow at bottom
add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "← Reviews · earnings · reputation flow back to all three sides ←",
         size=14, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

add_footer(s, 3)


# ─── SLIDE 4 — Market ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "Market",
            "₹6,400 Cr TAM. ₹900 Cr SOM in 5 years.",
            "Bigger than Indian food delivery was in 2013.")

# Left: stats table
table_top = Inches(2.7)
rows = [
    ("Gated communities (top 20 cities)", "28,000"),
    ("Corporate campuses / IT parks", "4,500"),
    ("Events per venue per year", "4 – 24"),
    ("Stall fee per event", "₹3K – ₹20K"),
    ("Total annual stall-booking GMV", "₹6,400 Cr"),
    ("Platform + BNPL take-rate", "12 – 18%"),
    ("Serviceable Obtainable Market (5-yr)", "₹900 Cr / yr"),
]
y = table_top
for i, (label, val) in enumerate(rows):
    highlight = i >= len(rows) - 2
    bg = AMBER_LIGHT if highlight else (GRAY_50 if i % 2 == 0 else WHITE)
    add_rect(s, Inches(0.6), y, Inches(7.5), Inches(0.48), bg)
    add_text(s, Inches(0.8), y + Inches(0.08), Inches(4.5), Inches(0.35),
             label, size=14, color=GRAY_700, bold=highlight)
    add_text(s, Inches(5.5), y + Inches(0.08), Inches(2.5), Inches(0.35),
             val, size=15, bold=True,
             color=AMBER if highlight else INDIGO, align=PP_ALIGN.RIGHT)
    y += Inches(0.5)

# Right: callout — Hyderabad beach-head
add_rect(s, Inches(8.5), Inches(2.7), Inches(4.3), Inches(3.8), INDIGO)
add_text(s, Inches(8.8), Inches(2.95), Inches(3.7), Inches(0.4),
         "HYDERABAD BEACH-HEAD", size=12, bold=True, color=AMBER)
add_text(s, Inches(8.8), Inches(3.5), Inches(3.7), Inches(1.2),
         "600", size=72, bold=True, color=WHITE)
add_text(s, Inches(8.8), Inches(4.7), Inches(3.7), Inches(0.5),
         "gated communities", size=16, color=WHITE)
add_text(s, Inches(8.8), Inches(5.2), Inches(3.7), Inches(0.5),
         "₹40 Cr annual GMV", size=20, bold=True, color=WHITE)
add_text(s, Inches(8.8), Inches(5.7), Inches(3.7), Inches(0.5),
         "→ 2% national market → our launchpad", size=12, color=AMBER_LIGHT)

add_footer(s, 4)


# ─── SLIDE 5 — Product (How it works) ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "Product",
            "The demand loop in 4 touch-points — already live.",
            "Every click creates data nobody else has.")

steps = [
    ("1", "RWA shares poll link",
     "One-tap share from venue dashboard\ninto resident WhatsApp group."),
    ("2", "Residents vote categories",
     "Food 40% · Fashion 25% · Plants 15%\nStandardised, deduped, rate-limited."),
    ("3", "Organiser creates event",
     "Demand chips pre-fill stallCategories.\nCopy-paste from past events."),
    ("4", "Vendor books with confidence",
     "Demand bars on event page.\nBNPL financing. Portable reputation."),
]
step_top = Inches(2.7)
step_h = Inches(3.8)
step_w = Inches(3.0)
x = Inches(0.6)
for num, title, body in steps:
    add_rect(s, x, step_top, step_w, step_h, WHITE, GRAY_300)
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3),
                                step_top + Inches(0.3), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = INDIGO
    circle.line.fill.background()
    add_text(s, x + Inches(0.3), step_top + Inches(0.3), Inches(0.7),
             Inches(0.7), num, size=24, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), step_top + Inches(1.3), step_w - Inches(0.6),
             Inches(0.6), title, size=16, bold=True, color=GRAY_900)
    add_text(s, x + Inches(0.3), step_top + Inches(2), step_w - Inches(0.6),
             Inches(1.6), body, size=12, color=GRAY_700)
    x += step_w + Inches(0.12)

add_footer(s, 5)


# ─── SLIDE 6 — Business Model ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "Business Model",
            "6 revenue streams. 7-week venue payback.",
            "We monetise the whole event lifecycle, not just listings.")

streams = [
    ("Booking commission", "5 – 8%", "on every stall fee"),
    ("BNPL interest", "2 – 3% / mo", "on pay-after-event stalls"),
    ("Setup-kit sales", "10 – 20%", "margin on tents + signage"),
    ("Smart-Score analytics", "₹2,500 / mo", "per organiser"),
    ("Premium venue listings", "₹1,000 / mo", "per venue"),
    ("Sponsor matchmaking", "10% cut", "on facilitated deals"),
]
grid_top = Inches(2.7)
grid_h = Inches(1.4)
grid_w = Inches(4.0)
for idx, (name, rate, note) in enumerate(streams):
    col = idx % 3
    row = idx // 3
    x = Inches(0.6) + (grid_w + Inches(0.12)) * col
    y = grid_top + (grid_h + Inches(0.15)) * row
    add_rect(s, x, y, grid_w, grid_h, GRAY_50, GRAY_300)
    add_rect(s, x, y, Inches(0.1), grid_h, INDIGO)
    add_text(s, x + Inches(0.3), y + Inches(0.2), grid_w - Inches(0.6),
             Inches(0.4), name, size=14, bold=True, color=GRAY_900)
    add_text(s, x + Inches(0.3), y + Inches(0.6), grid_w - Inches(0.6),
             Inches(0.4), rate, size=20, bold=True, color=INDIGO)
    add_text(s, x + Inches(0.3), y + Inches(1.0), grid_w - Inches(0.6),
             Inches(0.35), note, size=11, color=GRAY_500)

# Bottom unit economics bar
add_rect(s, Inches(0.6), Inches(5.9), Inches(12.15), Inches(1.1), AMBER_LIGHT)
add_text(s, Inches(0.9), Inches(6.0), Inches(4), Inches(0.4),
         "YEAR-2 UNIT ECONOMICS", size=11, bold=True, color=GRAY_700)
add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6),
         "₹1.44 L margin / venue / yr   ·   CAC: ₹6,000   ·   Payback: 7 weeks",
         size=18, bold=True, color=GRAY_900)

add_footer(s, 6)


# ─── SLIDE 7 — Traction ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "Traction",
            "Built. Live. Not slideware.",
            "http://69.62.80.48/stallmate — log in, walk the full loop.")

# Big numbers row
nums = [
    ("172", "files", INDIGO),
    ("30", "Prisma models", INDIGO),
    ("156", "API endpoints", INDIGO),
    ("60", "page routes", INDIGO),
]
x = Inches(0.6)
for val, lbl, color in nums:
    w = Inches(3.0)
    add_rect(s, x, Inches(2.7), w, Inches(1.8), GRAY_50, GRAY_300)
    add_text(s, x, Inches(2.85), w, Inches(1.0),
             val, size=60, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.95), w, Inches(0.5),
             lbl, size=14, color=GRAY_700, align=PP_ALIGN.CENTER)
    x += w + Inches(0.12)

# Seeded activity
add_rect(s, Inches(0.6), Inches(4.8), Inches(12.15), Inches(1.9), INDIGO)
add_text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.4),
         "REALISTIC SEEDED ACTIVITY", size=12, bold=True, color=AMBER)
add_text(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.5),
         "3 flagship venues · 53 real Hyderabad events · 4 organiser companies",
         size=16, bold=True, color=WHITE)
add_text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.5),
         "11 past events · 55 bookings · 29 reviews · 11 BNPL settlements · 90 demand votes",
         size=16, color=WHITE)
add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.4),
         "Missing = GTM, not product.",
         size=14, bold=True, color=AMBER)

add_footer(s, 7)


# ─── SLIDE 8 — Competition ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "Competition",
            "Listing directories don't compound. Data does.",
            "After 12 months per venue, our predictor beats the RWA's gut call.")

# Table
header_bg = GRAY_900
rows = [
    ("Player", "Model", "What they miss"),
    ("BookMyStall", "Listing directory", "No demand signal, no BNPL, no reputation"),
    ("EventsHigh", "Consumer ticketing", "B2C-only, not vendor-centric"),
    ("WhatsApp groups", "Zero-friction discovery", "Zero data, zero trust, zero scale"),
    ("Local organisers", "Hyper-local relationships", "Single-city ceiling, can't cross venues"),
    ("Stallmate", "3-sided marketplace + data loop", "— nothing ; we close the loop —"),
]
table_top = Inches(2.7)
row_h = Inches(0.55)
cols = [Inches(3.0), Inches(3.5), Inches(5.7)]
col_x = [Inches(0.6), Inches(0.6) + cols[0], Inches(0.6) + cols[0] + cols[1]]
for i, row in enumerate(rows):
    y = table_top + row_h * i
    is_header = i == 0
    is_us = i == len(rows) - 1
    bg = header_bg if is_header else (AMBER_LIGHT if is_us else (GRAY_50 if i % 2 == 0 else WHITE))
    for cx, cw, text in zip(col_x, cols, row):
        add_rect(s, cx, y, cw, row_h, bg, GRAY_300)
        add_text(s, cx + Inches(0.15), y + Inches(0.15), cw - Inches(0.3),
                 row_h - Inches(0.2), text,
                 size=12 if not is_header else 11,
                 bold=is_header or is_us,
                 color=WHITE if is_header else (GRAY_900 if is_us else GRAY_700))

# Moat callout
add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
         "Our moat: resident demand graph + vendor reputation graph + event-outcome dataset",
         size=14, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

add_footer(s, 8)


# ─── SLIDE 9 — GTM & Roadmap ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
title_block(s, "Go-to-Market",
            "Hyderabad → 4 cities → Tier 2.",
            "Every city launch uses the playbook from the city before.")

phases = [
    ("Phase 1", "Next 90 days", "Hyderabad beach-head", INDIGO,
     ["15 pilot gated communities", "5 organisers live",
      "300 vendors onboarded", "120 events · ₹2.5 Cr GMV · ₹18 L rev"]),
    ("Phase 2", "Months 4–12", "4-city expansion", AMBER,
     ["Bangalore, Pune, Chennai launched",
      "BNPL live with NBFC partner",
      "₹5 Cr credit line deployed",
      "500 vendors / mo · ₹4 Cr ARR"]),
    ("Phase 3", "Year 2+", "Tier-2 + extensions", GREEN,
     ["Vizag, Indore, Coimbatore, Kochi",
      "Fundraisers, schools, festivals",
      "₹40 Cr annual GMV run-rate",
      "Series A ready"]),
]
phase_top = Inches(2.7)
phase_w = Inches(4.0)
x = Inches(0.6)
for p, when, name, color, bullets in phases:
    add_rect(s, x, phase_top, phase_w, Inches(4), WHITE, GRAY_300)
    add_rect(s, x, phase_top, phase_w, Inches(0.6), color)
    add_text(s, x + Inches(0.25), phase_top + Inches(0.12), phase_w - Inches(0.5),
             Inches(0.4), f"{p} · {when}", size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), phase_top + Inches(0.85), phase_w - Inches(0.5),
             Inches(0.5), name, size=18, bold=True, color=GRAY_900)
    y = phase_top + Inches(1.5)
    for b in bullets:
        add_text(s, x + Inches(0.25), y, phase_w - Inches(0.5), Inches(0.5),
                 "•  " + b, size=13, color=GRAY_700)
        y += Inches(0.55)
    x += phase_w + Inches(0.12)

add_footer(s, 9)


# ─── SLIDE 10 — Ask ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
# Indigo hero on left
add_rect(s, 0, 0, Inches(5.5), SH, INDIGO)
add_text(s, Inches(0.8), Inches(1.2), Inches(4.5), Inches(0.4),
         "THE ASK", size=14, bold=True, color=AMBER)
add_text(s, Inches(0.8), Inches(1.8), Inches(4.5), Inches(1.5),
         "₹8 Cr", size=110, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.8), Inches(4.5), Inches(0.5),
         "Seed round", size=22, color=WHITE)
add_text(s, Inches(0.8), Inches(4.8), Inches(4.5), Inches(0.5),
         "Close in 45 days.", size=16, color=AMBER)
add_text(s, Inches(0.8), Inches(5.3), Inches(4.5), Inches(0.5),
         "Deploy Hyderabad in 60.", size=16, color=AMBER)
add_text(s, Inches(0.8), Inches(6.4), Inches(4.5), Inches(0.3),
         "your@email · +91 XXXXX XXXXX", size=11, color=WHITE)

# Right: use of funds + milestones
title_block(s, "", "Use of Funds", top=Inches(0.8))
fund = [
    ("40%", "GTM", "venue & vendor acquisition", INDIGO),
    ("25%", "Product", "BNPL, mobile, ML prediction", AMBER),
    ("20%", "Operations", "3 city launches, inventory", GREEN),
    ("10%", "Compliance", "FSSAI, MSME, schemes", GRAY_500),
    ("5%", "Reserve", "", GRAY_300),
]
y = Inches(2.0)
for pct, label, detail, color in fund:
    add_text(s, Inches(6.3), y, Inches(1.2), Inches(0.5),
             pct, size=22, bold=True, color=color)
    add_text(s, Inches(7.6), y + Inches(0.05), Inches(2), Inches(0.4),
             label, size=14, bold=True, color=GRAY_900)
    add_text(s, Inches(7.6), y + Inches(0.4), Inches(5), Inches(0.4),
             detail, size=12, color=GRAY_500)
    y += Inches(0.55)

# Milestones strip
add_rect(s, Inches(6), Inches(5.4), Inches(6.8), Inches(1.5), AMBER_LIGHT)
add_text(s, Inches(6.3), Inches(5.5), Inches(6.3), Inches(0.4),
         "WHAT YOU GET", size=11, bold=True, color=GRAY_700)
add_text(s, Inches(6.3), Inches(5.85), Inches(6.3), Inches(0.4),
         "M3: 15 paying pilot venues", size=12, color=GRAY_900)
add_text(s, Inches(6.3), Inches(6.2), Inches(6.3), Inches(0.4),
         "M9: 500 vendors transacting / mo", size=12, color=GRAY_900)
add_text(s, Inches(6.3), Inches(6.55), Inches(6.3), Inches(0.4),
         "M18: Series A ready, ₹4 Cr ARR", size=12, color=GRAY_900)

add_footer(s, 10)


# ── Save ──────────────────────────────────────────────────────────────────
import os
out = "/root/stallmate/public/stallmate-deck.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"Saved: {out}")
